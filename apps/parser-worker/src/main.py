from __future__ import annotations

import asyncio
import base64
import html
import logging
import os
import re
import subprocess
import tempfile
import time
import uuid
from contextlib import asynccontextmanager
from pathlib import Path
from typing import Any

from fastapi import BackgroundTasks, Depends, FastAPI, File, Form, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from pydantic import BaseModel

logger = logging.getLogger('parser-worker')
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(name)s: %(message)s")

# Clean proxy environment for httpx/huggingface_hub compatibility
for k in ["ALL_PROXY", "all_proxy"]:
    if os.environ.get(k, "").startswith("socks://"):
        os.environ.pop(k, None)

UPLOAD_ROOT = Path(os.environ.get("UPLOAD_ROOT", "/tmp/llmwiki/parser"))
MAX_FILE_BYTES = 200 * 1024 * 1024
SUPPORTED_EXTENSIONS = {".md", ".txt", ".csv", ".html", ".htm", ".doc", ".docx", ".pdf", ".xls", ".xlsx", ".pptx", ".png", ".jpg", ".jpeg"}
ANTIWORD_BIN = os.environ.get("ANTIWORD_BIN", "antiword")
DOCLING_TIMEOUT_SECONDS = float(os.environ.get("DOCLING_TIMEOUT_SECONDS", "240"))
PDF_PARSE_MODE = os.environ.get("PDF_PARSE_MODE", "hybrid").lower()
OCR_PROVIDER = os.environ.get("OCR_PROVIDER", "none").lower()
OCR_TIMEOUT_SECONDS = float(os.environ.get("OCR_TIMEOUT_SECONDS", "900"))
OCR_POLL_INTERVAL_SECONDS = float(os.environ.get("OCR_POLL_INTERVAL_SECONDS", "5"))
OCR_MAX_FILE_BYTES = int(os.environ.get("OCR_MAX_FILE_BYTES", str(50 * 1024 * 1024)))
BAIDU_OCR_API_KEY = os.environ.get("BAIDU_OCR_API_KEY", "").strip()
BAIDU_OCR_SECRET_KEY = os.environ.get("BAIDU_OCR_SECRET_KEY", "").strip()
BAIDU_OCR_ENDPOINT = os.environ.get(
    "BAIDU_OCR_ENDPOINT", "https://aip.baidubce.com"
).rstrip("/")
LOCAL_DOCLING_ENABLED = os.environ.get("LOCAL_DOCLING_ENABLED", "1").lower() not in {
    "0",
    "false",
    "no",
}
tasks: dict[str, dict[str, Any]] = {}
_torchvision_compat_lib = None
_baidu_access_token: tuple[str, float] | None = None


def _pdf_native_quality(text: str) -> str:
    if not text:
        return "empty"
    replacement_ratio = text.count("\ufffd") / max(len(text), 1)
    control_count = sum(1 for char in text if ord(char) < 32 and char not in "\n\r\t")
    control_ratio = control_count / max(len(text), 1)
    return "poor" if replacement_ratio > 0.08 or control_ratio > 0.02 else "good"


def classify_pdf(
    page_count: int, text_pages: int, native_chars: int, native_quality: str
) -> str:
    """Classify a PDF from page coverage, not just total extracted length."""
    page_ratio = text_pages / max(page_count, 1)
    average_chars = native_chars / max(page_count, 1)
    if (
        native_quality == "good"
        and native_chars >= 80
        and (page_count == 1 or page_ratio >= 0.65)
        and average_chars >= 80
    ):
        return "text"
    if page_ratio <= 0.20 or (native_chars < 40 and native_quality != "good"):
        return "scanned"
    return "mixed"

async def periodic_cleanup():
    while True:
        await asyncio.sleep(300)
        current_time = time.time()
        for tid in list(tasks.keys()):
            t_info = tasks[tid]
            if t_info.get("status") in ("completed", "failed"):
                if current_time - t_info.get("created_at", current_time) > 1800:
                    del tasks[tid]

@asynccontextmanager
async def lifespan(app: FastAPI):
    UPLOAD_ROOT.mkdir(parents=True, exist_ok=True)
    cleanup_task = asyncio.create_task(periodic_cleanup())
    yield
    cleanup_task.cancel()

app = FastAPI(title="LLMWiki Parser Worker", version="0.4.0", lifespan=lifespan)

allowed_origins = os.environ.get('CORS_ORIGINS', 'http://localhost:3000,http://localhost:3001').split(',')
app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

security = HTTPBearer(auto_error=False)

def verify_auth(credentials: HTTPAuthorizationCredentials = Depends(security)):
    token = os.environ.get("AUTH_TOKEN")
    if token:
        if not credentials or credentials.credentials != token:
            raise HTTPException(status_code=401, detail="Invalid or missing authentication token")

class ParseResponse(BaseModel):
    task_id: str
    status: str
    message: str

@app.get("/health")
def health_check():
    return {
        "status": "healthy",
        "version": app.version,
        "pdf_parse_mode": PDF_PARSE_MODE,
        "ocr_provider": OCR_PROVIDER,
        "local_docling_enabled": LOCAL_DOCLING_ENABLED,
    }

@app.get("/metrics")
def metrics():
    total = len(tasks)
    by_status = {}
    by_engine = {}
    by_classification = {}
    for t in tasks.values():
        s = t.get('status', 'unknown')
        by_status[s] = by_status.get(s, 0) + 1
        if t.get("engine"):
            engine = str(t["engine"])
            by_engine[engine] = by_engine.get(engine, 0) + 1
        if t.get("classification"):
            classification = str(t["classification"])
            by_classification[classification] = by_classification.get(classification, 0) + 1
    return {
        'total_tasks': total,
        'by_status': by_status,
        'by_engine': by_engine,
        'by_classification': by_classification,
    }

def extract_plaintext(filename: str, content: bytes) -> str:
    suffix = Path(filename).suffix.lower()
    text = content.decode("utf-8-sig", errors="replace")
    if suffix in {".html", ".htm"}:
        text = html.unescape(re.sub(r"<[^>]+>", " ", text))
    return text.strip()


def normalize_markdown(markdown: str, filename: str) -> str:
    """Normalize parser output without flattening meaningful document structure.

    Office converters commonly emit a title once as metadata and once as body
    text, and legacy Word emits form-feed page breaks. Removing only repeated
    standalone title lines keeps the original wording while preventing the
    duplicate title from becoming a second high-ranking retrieval passage.
    """
    title = Path(filename).stem.strip()
    lines = markdown.replace("\r\n", "\n").replace("\r", "\n").replace("\x0c", "\n\n").split("\n")
    normalized: list[str] = []
    title_seen = False
    for raw_line in lines:
        line = raw_line.replace("\u200b", "").replace("\ufeff", "").replace("\xa0", " ").rstrip()
        comparable = re.sub(r"^\s*#+\s*", "", line).strip()
        if title and comparable == title:
            if title_seen:
                continue
            title_seen = True
        normalized.append(line)

    result = "\n".join(normalized)
    result = re.sub(r"\n{3,}", "\n\n", result).strip()
    result = re.sub(r"([^\n])\n(第\s*[\d一二三四五六七八九十百千万〇零两]+\s*[章节条款项])", r"\1\n\n\2", result)
    return result

def extract_legacy_word(path: Path) -> str:
    env = os.environ.copy()
    try:
        result = subprocess.run(
            [ANTIWORD_BIN, "-f", str(path)],
            check=False,
            capture_output=True,
            timeout=120,
            env=env,
        )
    except FileNotFoundError as exc:
        raise RuntimeError("Legacy .doc conversion is unavailable: antiword is not installed") from exc
    except subprocess.TimeoutExpired as exc:
        raise RuntimeError("Legacy .doc conversion timed out after 120 seconds") from exc
    if result.returncode != 0:
        detail = result.stderr.decode("utf-8", errors="replace").strip()
        raise RuntimeError(f"Legacy .doc conversion failed{f': {detail}' if detail else ''}")
    text = result.stdout.decode("utf-8", errors="replace").strip()
    if not text:
        raise RuntimeError("Legacy .doc conversion returned empty text")
    return text

def extract_docx(path: Path) -> str:
    """Extract .docx to Markdown while preserving paragraph/table order."""
    try:
        import docx
        from docx.document import Document as DocxDocument
        from docx.table import Table
        from docx.text.paragraph import Paragraph
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P

        doc = docx.Document(str(path))
        lines = []

        def iter_blocks(parent: DocxDocument):
            for child in parent.element.body.iterchildren():
                if isinstance(child, CT_P):
                    yield Paragraph(child, parent)
                elif isinstance(child, CT_Tbl):
                    yield Table(child, parent)

        for block in iter_blocks(doc):
            if isinstance(block, Paragraph):
                txt = block.text.strip()
                if not txt:
                    continue
                style_name = (block.style.name if block.style else "").lower()
                heading_match = re.search(r"heading\s*([1-6])", style_name)
                lines.append(f"{'#' * int(heading_match.group(1))} {txt}" if heading_match else txt)
            else:
                rows = []
                for row in block.rows:
                    cells = [cell.text.strip().replace("\n", " ").replace("|", "\\|") for cell in row.cells]
                    if any(cells):
                        rows.append(cells)
                if rows:
                    width = max(len(row) for row in rows)
                    normalized = [row + [""] * (width - len(row)) for row in rows]
                    t_lines = ["| " + " | ".join(normalized[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
                    t_lines.extend("| " + " | ".join(row) + " |" for row in normalized[1:])
                    lines.append("\n".join(t_lines))
        return "\n\n".join(lines).strip()
    except Exception as e:
        logger.warning(f"python-docx extraction failed for {path}: {e}")
        return ""

def inspect_pdf_native(path: Path) -> dict[str, Any]:
    """Inspect native PDF text page-by-page without OCR or GPU work."""
    result: dict[str, Any] = {
        "markdown": "",
        "page_count": 0,
        "text_pages": 0,
        "native_chars": 0,
        "native_page_ratio": 0.0,
        "native_quality": "empty",
        "classification": "unknown",
        "page_texts": [],
        "native_page_indexes": [],
    }
    try:
        import pypdf

        reader = pypdf.PdfReader(str(path))
        result["page_count"] = len(reader.pages)
        pages_text = []
        page_texts = []
        for i, page in enumerate(reader.pages):
            txt = (page.extract_text() or "").strip()
            page_texts.append(txt)
            result["native_chars"] += len(re.sub(r"\s+", "", txt))
            if len(re.sub(r"\s+", "", txt)) >= 40:
                result["text_pages"] += 1
            if txt:
                pages_text.append(f"## 第 {i+1} 页\n\n{txt}")
        if pages_text:
            result["markdown"] = f"# {path.stem}\n\n" + "\n\n---\n\n".join(pages_text)
        result["page_texts"] = page_texts
        result["native_page_ratio"] = result["text_pages"] / max(result["page_count"], 1)
        result["native_quality"] = _pdf_native_quality(result["markdown"])
        result["classification"] = classify_pdf(
            int(result["page_count"]),
            int(result["text_pages"]),
            int(result["native_chars"]),
            str(result["native_quality"]),
        )
        if result["native_quality"] == "good":
            result["native_page_indexes"] = [
                i
                for i, txt in enumerate(page_texts)
                if len(re.sub(r"\s+", "", txt)) >= 40
            ]
    except Exception as e:
        logger.warning(f"pypdf extraction failed for {path}: {e}")
        result["error"] = str(e)
    return result


def create_pdf_subset(path: Path, page_indexes: list[int]) -> Path:
    """Create a short-lived PDF containing only selected pages."""
    import pypdf

    reader = pypdf.PdfReader(str(path))
    handle = tempfile.NamedTemporaryFile(
        prefix="ocr-pages-", suffix=".pdf", dir=str(UPLOAD_ROOT), delete=False
    )
    subset_path = Path(handle.name)
    try:
        writer = pypdf.PdfWriter()
        for index in page_indexes:
            if 0 <= index < len(reader.pages):
                writer.add_page(reader.pages[index])
        with handle:
            writer.write(handle)
        if not page_indexes:
            raise RuntimeError("No scan pages selected for OCR")
        return subset_path
    except Exception:
        handle.close()
        subset_path.unlink(missing_ok=True)
        raise


def split_ocr_pages(markdown: str) -> list[str]:
    """Split provider Markdown into page bodies when page headings exist."""
    matches = list(re.finditer(r"(?m)^##\s*(?:第\s*)?\d+\s*页\s*$", markdown))
    if not matches:
        return [markdown.strip()] if markdown.strip() else []
    pages = []
    for index, match in enumerate(matches):
        start = match.end()
        end = matches[index + 1].start() if index + 1 < len(matches) else len(markdown)
        body = re.sub(r"\n\s*---\s*$", "", markdown[start:end]).strip()
        if body:
            pages.append(body)
    return pages


def merge_mixed_pdf_markdown(
    title: str,
    page_texts: list[str],
    scan_page_indexes: list[int],
    ocr_markdown: str,
) -> str:
    """Put native and OCR page bodies back into the original page order."""
    ocr_pages = split_ocr_pages(ocr_markdown)
    if not ocr_pages:
        raise RuntimeError("OCR returned no page content for mixed PDF")
    if len(ocr_pages) != len(scan_page_indexes):
        logger.warning(
            "OCR page count mismatch for mixed PDF: expected=%s actual=%s; "
            "assigning provider output to the first scan page",
            len(scan_page_indexes),
            len(ocr_pages),
        )
        ocr_pages = ["\n\n".join(ocr_pages)]

    ocr_by_page = {
        page_index: ocr_pages[position]
        for position, page_index in enumerate(scan_page_indexes[: len(ocr_pages)])
    }
    merged_pages = []
    for page_index in range(max(len(page_texts), max(scan_page_indexes, default=-1) + 1)):
        body = ocr_by_page.get(page_index) or (
            page_texts[page_index] if page_index < len(page_texts) else ""
        )
        if body.strip():
            merged_pages.append(f"## 第 {page_index + 1} 页\n\n{body.strip()}")
    if not merged_pages:
        raise RuntimeError("Mixed PDF merge produced empty Markdown")
    return f"# {title}\n\n" + "\n\n---\n\n".join(merged_pages)


def extract_pdf_native(path: Path) -> str:
    """Extract native PDF text & structure via pypdf in milliseconds."""
    return str(inspect_pdf_native(path).get("markdown", ""))

def extract_excel(path: Path) -> str:
    """Extract .xlsx / .xls sheets to Markdown tables."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True)
        sheets_md = []
        for sheetname in wb.sheetnames:
            sheet = wb[sheetname]
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue
            table_lines = [f"### 工作表：{sheetname}\n"]
            header = [str(cell if cell is not None else "") for cell in rows[0]]
            table_lines.append("| " + " | ".join(header) + " |")
            table_lines.append("| " + " | ".join(["---"] * len(header)) + " |")
            for row in rows[1:]:
                if all(c is None or str(c).strip() == "" for c in row):
                    continue
                cells = [str(c if c is not None else "").replace("\n", " ") for c in row]
                table_lines.append("| " + " | ".join(cells) + " |")
            sheets_md.append("\n".join(table_lines))
        return "\n\n".join(sheets_md)
    except Exception as openpyxl_error:
        # openpyxl intentionally does not read the legacy BIFF .xls format.
        # Keep the lightweight xlrd route optional so production can support
        # .xls without installing the heavyweight layout engine.
        logger.info(f"openpyxl extraction unavailable for {path}: {openpyxl_error}")
        try:
            import xlrd

            workbook = xlrd.open_workbook(str(path), on_demand=True)
            sheets_md = []
            for sheet in workbook.sheets():
                if sheet.nrows == 0:
                    continue
                rows = [
                    [str(sheet.cell_value(row, col) or "").replace("|", "\\|").replace("\n", " ") for col in range(sheet.ncols)]
                    for row in range(sheet.nrows)
                ]
                while rows and not any(rows[-1]):
                    rows.pop()
                if not rows:
                    continue
                width = max(len(row) for row in rows)
                rows = [row + [""] * (width - len(row)) for row in rows]
                lines = [f"### 工作表：{sheet.name}\n"]
                lines.append("| " + " | ".join(rows[0]) + " |")
                lines.append("| " + " | ".join(["---"] * width) + " |")
                lines.extend("| " + " | ".join(row) + " |" for row in rows[1:])
                sheets_md.append("\n".join(lines))
            return "\n\n".join(sheets_md)
        except Exception as xlrd_error:
            logger.warning(f"Legacy Excel extraction failed for {path}: {xlrd_error}")
            return ""


def _pptx_table_markdown(table: Any) -> str:
    rows = []
    for row in table.rows:
        cells = [
            str(cell.text or "").strip().replace("|", "\\|").replace("\n", " ")
            for cell in row.cells
        ]
        if any(cells):
            rows.append(cells)
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    rows = [row + [""] * (width - len(row)) for row in rows]
    lines = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in rows[1:])
    return "\n".join(lines)


def extract_pptx_native(path: Path) -> tuple[list[str], list[dict[str, Any]]]:
    """Extract slide text/tables and retain embedded images for OCR.

    This is the production fallback when local Docling is intentionally
    disabled. Native text is never discarded just because a slide also has a
    picture; pictures are OCR'ed separately when a cloud provider is configured.
    """
    from pptx import Presentation

    presentation = Presentation(str(path))
    slide_blocks: list[str] = []
    image_parts: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape_number, shape in enumerate(slide.shapes, start=1):
            if getattr(shape, "has_text_frame", False):
                text = "\n".join(
                    paragraph.text.strip()
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                ).strip()
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False):
                table_md = _pptx_table_markdown(shape.table)
                if table_md:
                    parts.append(table_md)
            if getattr(shape, "shape_type", None) == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image = shape.image
                    image_parts.append({
                        "slide": slide_number,
                        "shape": shape_number,
                        "ext": str(image.ext or "png"),
                        "blob": image.blob,
                    })
                except Exception as image_error:
                    logger.warning("Unable to extract PPTX image on slide %s: %s", slide_number, image_error)
        slide_blocks.append("\n\n".join(parts).strip())
    return slide_blocks, image_parts

async def convert_with_docling(path: Path) -> str:
    """Docling deep layout extraction with compatibility guard."""
    def _run():
        for k in ["ALL_PROXY", "all_proxy"]:
            if os.environ.get(k, "").startswith("socks://"):
                os.environ.pop(k, None)
        global _torchvision_compat_lib
        import torch
        if _torchvision_compat_lib is None:
            try:
                _torchvision_compat_lib = torch.library.Library("torchvision", "DEF")
            except RuntimeError:
                _torchvision_compat_lib = torch.library.Library("torchvision", "FRAGMENT")
        for operator in ("nms", "qnms"):
            try:
                _torchvision_compat_lib.define(
                    f"{operator}(Tensor boxes, Tensor scores, float iou_threshold) -> Tensor"
                )
            except Exception:
                pass
        from docling.document_converter import DocumentConverter
        converter = DocumentConverter()
        result = converter.convert(str(path))
        return result.document.export_to_markdown()
    return await asyncio.to_thread(_run)


async def _baidu_token(client: Any, api_key: str, secret_key: str, endpoint: str) -> str:
    global _baidu_access_token
    now = time.time()
    if _baidu_access_token and _baidu_access_token[1] > now + 60:
        return _baidu_access_token[0]
    if not api_key or not secret_key:
        raise RuntimeError("Baidu OCR is enabled but API key/secret key is not configured")
    response = await client.post(
        f"{endpoint}/oauth/2.0/token",
        data={
            "grant_type": "client_credentials",
            "client_id": api_key,
            "client_secret": secret_key,
        },
    )
    response.raise_for_status()
    payload = response.json()
    token = str(payload.get("access_token") or "")
    if not token:
        raise RuntimeError(f"Baidu OCR token request failed: {payload.get('error_description') or payload}")
    expires_in = int(payload.get("expires_in") or 2592000)
    _baidu_access_token = (token, now + max(expires_in, 300))
    return token


async def convert_with_baidu_ocr(
    path: Path, ocr_config: dict[str, str]
) -> tuple[str, dict[str, Any]]:
    """Use Baidu's async document parser for scanned/mixed PDFs.

    This is intentionally a PDF-level API call: it preserves page boundaries,
    headings and tables better than rendering every page and calling generic OCR.
    """
    if path.stat().st_size > OCR_MAX_FILE_BYTES:
        raise RuntimeError(
            f"PDF is {path.stat().st_size} bytes; Baidu file_data limit is {OCR_MAX_FILE_BYTES} bytes"
        )
    try:
        import httpx
    except ImportError as exc:
        raise RuntimeError("Baidu OCR requires the httpx dependency") from exc

    timeout = httpx.Timeout(OCR_TIMEOUT_SECONDS, connect=30.0)
    async with httpx.AsyncClient(timeout=timeout) as client:
        endpoint_base = str(ocr_config.get("endpoint") or BAIDU_OCR_ENDPOINT).rstrip("/")
        token = await _baidu_token(
            client,
            str(ocr_config.get("api_key") or BAIDU_OCR_API_KEY),
            str(ocr_config.get("secret_key") or BAIDU_OCR_SECRET_KEY),
            endpoint_base,
        )
        endpoint = f"{endpoint_base}/rest/2.0/brain/online/v2/parser/task"
        raw = await asyncio.to_thread(path.read_bytes)
        response = await client.post(
            endpoint,
            params={"access_token": token},
            data={
                "file_data": base64.b64encode(raw).decode("ascii"),
                "file_name": path.name,
                "language_type": "CHN_ENG",
                "angle_adjust": "true",
                "html_table_format": "false",
            },
            headers={"Content-Type": "application/x-www-form-urlencoded"},
        )
        response.raise_for_status()
        submitted = response.json()
        if int(submitted.get("error_code", 0) or 0) != 0:
            raise RuntimeError(
                f"Baidu OCR submit failed: {submitted.get('error_msg') or submitted}"
            )
        task_id = str((submitted.get("result") or {}).get("task_id") or "")
        if not task_id:
            raise RuntimeError(f"Baidu OCR did not return task_id: {submitted}")

        query_endpoint = f"{endpoint_base}/rest/2.0/brain/online/v2/parser/task/query"
        deadline = time.monotonic() + OCR_TIMEOUT_SECONDS
        status_payload: dict[str, Any] = {}
        while time.monotonic() < deadline:
            await asyncio.sleep(OCR_POLL_INTERVAL_SECONDS)
            status_response = await client.post(
                query_endpoint,
                params={"access_token": token},
                data={"task_id": task_id},
                headers={"Content-Type": "application/x-www-form-urlencoded"},
            )
            status_response.raise_for_status()
            status_payload = status_response.json()
            detail = status_payload.get("result") or {}
            status = str(detail.get("status") or "")
            if status == "success":
                markdown_url = str(detail.get("markdown_url") or "")
                if not markdown_url:
                    raise RuntimeError(f"Baidu OCR returned no markdown URL: {status_payload}")
                markdown_response = await client.get(markdown_url)
                markdown_response.raise_for_status()
                markdown = markdown_response.text.strip()
                if not markdown:
                    raise RuntimeError("Baidu OCR returned empty Markdown")
                return markdown, {
                    "ocr_provider": "baidu",
                    "ocr_task_id": task_id,
                    "ocr_cost_pages": detail.get("cost_page_num"),
                }
            if status == "failed":
                raise RuntimeError(
                    f"Baidu OCR task failed: {detail.get('task_error') or status_payload}"
                )
        raise TimeoutError(f"Baidu OCR task timed out after {OCR_TIMEOUT_SECONDS:g} seconds")


async def convert_image_with_baidu_ocr(
    path: Path, ocr_config: dict[str, str]
) -> tuple[str, dict[str, Any]]:
    """OCR a standalone image or an embedded PPTX image with Baidu.

    The document parser endpoint is the preferred route for PDF because it
    preserves page structure. Images have no page/document container, so use
    Baidu's high-accuracy OCR endpoint and keep the returned line order.
    """
    provider = str(ocr_config.get("provider") or OCR_PROVIDER).lower()
    if provider != "baidu":
        raise RuntimeError("Image OCR requires a Baidu OCR provider or local Docling")
    if path.stat().st_size > OCR_MAX_FILE_BYTES:
        raise RuntimeError(
            f"Image is {path.stat().st_size} bytes; OCR file limit is {OCR_MAX_FILE_BYTES} bytes"
        )
    try:
        import httpx
    except ImportError as exc:
        raise RuntimeError("Baidu OCR requires the httpx dependency") from exc

    timeout = httpx.Timeout(OCR_TIMEOUT_SECONDS, connect=30.0)
    async with httpx.AsyncClient(timeout=timeout) as client:
        endpoint_base = str(ocr_config.get("endpoint") or BAIDU_OCR_ENDPOINT).rstrip("/")
        token = await _baidu_token(
            client,
            str(ocr_config.get("api_key") or BAIDU_OCR_API_KEY),
            str(ocr_config.get("secret_key") or BAIDU_OCR_SECRET_KEY),
            endpoint_base,
        )
        raw = await asyncio.to_thread(path.read_bytes)
        response = await client.post(
            f"{endpoint_base}/rest/2.0/ocr/v1/accurate_basic",
            params={"access_token": token},
            data={
                "image": base64.b64encode(raw).decode("ascii"),
                "detect_direction": "true",
                "paragraph": "true",
                "probability": "true",
            },
            headers={"Content-Type": "application/x-www-form-urlencoded"},
        )
        response.raise_for_status()
        payload = response.json()
        if int(payload.get("error_code", 0) or 0) != 0:
            raise RuntimeError(
                f"Baidu image OCR failed: {payload.get('error_msg') or payload}"
            )
        words_result = payload.get("words_result") or []
        lines: list[str] = []
        probabilities: list[float] = []
        for item in words_result:
            words = str(item.get("words") or "").strip() if isinstance(item, dict) else str(item).strip()
            if not words:
                continue
            lines.append(words)
            if isinstance(item, dict):
                probability = item.get("probability")
                if isinstance(probability, dict):
                    probability = probability.get("average")
                try:
                    if probability is not None:
                        probabilities.append(float(probability))
                except (TypeError, ValueError):
                    pass
        if not lines:
            raise RuntimeError("Baidu image OCR returned no text")
        metadata: dict[str, Any] = {
            "ocr_provider": "baidu",
            "ocr_endpoint": "accurate_basic",
            "ocr_words_result_num": len(lines),
        }
        if probabilities:
            metadata["ocr_average_confidence"] = round(sum(probabilities) / len(probabilities), 4)
        return "\n".join(lines), metadata


async def convert_with_cloud_ocr(
    path: Path, ocr_config: dict[str, str]
) -> tuple[str, dict[str, Any]]:
    provider = str(ocr_config.get("provider") or OCR_PROVIDER).lower()
    if provider == "baidu":
        return await convert_with_baidu_ocr(path, ocr_config)
    raise RuntimeError(
        "Scanned PDF requires an OCR provider configured as baidu (or local Docling)"
    )


async def convert_pptx_without_docling(
    path: Path, ocr_config: dict[str, str]
) -> tuple[str, str, dict[str, Any]]:
    """Build a slide-preserving Markdown representation without Docling.

    Production deliberately does not install/enable Docling. Native PPTX
    text and tables remain lossless; embedded images are sent through the
    configured OCR provider. If a slide contains images and no OCR route is
    available, fail closed instead of silently indexing incomplete content.
    """
    slide_blocks, image_parts = await asyncio.to_thread(extract_pptx_native, path)
    provider = str(ocr_config.get("provider") or OCR_PROVIDER).lower()
    image_by_slide: dict[int, list[str]] = {}
    metadata: dict[str, Any] = {
        "slide_count": len(slide_blocks),
        "embedded_image_count": len(image_parts),
    }
    if image_parts and provider == "none":
        raise RuntimeError(
            "PPTX contains embedded images; configure Baidu OCR for complete image text extraction"
        )
    for position, image in enumerate(image_parts, start=1):
        temp = tempfile.NamedTemporaryFile(
            prefix=f"pptx-image-{position}-", suffix=f".{image['ext']}", dir=str(UPLOAD_ROOT), delete=False
        )
        image_path = Path(temp.name)
        try:
            with temp:
                temp.write(image["blob"])
            image_md, image_metadata = await convert_image_with_baidu_ocr(image_path, ocr_config)
            image_by_slide.setdefault(int(image["slide"]), []).append(
                f"### 图片区域 {image['shape']}\n\n{image_md.strip()}"
            )
            for key, value in image_metadata.items():
                if key == "ocr_average_confidence" and value is not None:
                    previous = metadata.get(key)
                    metadata[key] = round(
                        (float(previous) * (position - 1) + float(value)) / position, 4
                    ) if previous is not None else value
                elif key.startswith("ocr_"):
                    metadata[key] = metadata.get(key, 0) + value if isinstance(value, (int, float)) else value
        finally:
            image_path.unlink(missing_ok=True)

    sections = [f"# {path.stem}"]
    for slide_number, block in enumerate(slide_blocks, start=1):
        content = [block] if block else []
        content.extend(image_by_slide.get(slide_number, []))
        if content:
            sections.append(f"## 第 {slide_number} 页\n\n" + "\n\n".join(content))
    markdown = "\n\n---\n\n".join(sections)
    if not any(block.strip() for block in slide_blocks) and not image_by_slide:
        raise RuntimeError("PPTX contains no extractable text, tables, or OCR results")
    return markdown, "python-pptx-native+ocr" if image_parts else "python-pptx-native", metadata


async def convert_pdf_with_fallback(
    path: Path,
    classification: str,
    native_md: str,
    ocr_config: dict[str, str],
    page_texts: list[str] | None = None,
    native_page_indexes: list[int] | None = None,
) -> tuple[str, str, dict[str, Any]]:
    """Route PDF to the cheapest suitable engine, then fail open safely."""
    metadata: dict[str, Any] = {}
    if classification == "text" and PDF_PARSE_MODE in {"fast", "hybrid", "auto"}:
        return native_md, "pypdf-native", metadata

    if classification in {"scanned", "mixed"} and str(
        ocr_config.get("provider") or OCR_PROVIDER
    ).lower() != "none":
        ocr_path = path
        ocr_subset_path: Path | None = None
        try:
            scan_page_indexes = [
                index
                for index in range(len(page_texts or []))
                if index not in set(native_page_indexes or [])
            ]
            if (
                classification == "mixed"
                and page_texts
                and scan_page_indexes
                and len(scan_page_indexes) < len(page_texts)
            ):
                ocr_subset_path = await asyncio.to_thread(
                    create_pdf_subset, path, scan_page_indexes
                )
                ocr_path = ocr_subset_path
            markdown, ocr_metadata = await convert_with_cloud_ocr(ocr_path, ocr_config)
            if ocr_subset_path:
                markdown = merge_mixed_pdf_markdown(
                    path.stem,
                    page_texts or [],
                    scan_page_indexes,
                    markdown,
                )
                ocr_metadata = {
                    **ocr_metadata,
                    "ocr_original_pages": [index + 1 for index in scan_page_indexes],
                    "ocr_cost_pages": len(scan_page_indexes),
                }
                return markdown, "ocr-baidu-mixed-pages", ocr_metadata
            return markdown, f"ocr-{ocr_config.get('provider') or OCR_PROVIDER}", ocr_metadata
        except Exception as ocr_err:
            logger.warning(f"Cloud OCR on {path.name} failed: {ocr_err}")
            metadata["ocr_error"] = str(ocr_err)
        finally:
            if ocr_subset_path:
                ocr_subset_path.unlink(missing_ok=True)

    if LOCAL_DOCLING_ENABLED:
        try:
            markdown = await asyncio.wait_for(
                convert_with_docling(path), timeout=DOCLING_TIMEOUT_SECONDS
            )
            return markdown, "docling-local", metadata
        except Exception as docling_err:
            logger.warning(f"Local Docling on {path.name} failed/timed out: {docling_err}")
            metadata["docling_error"] = str(docling_err)

    if native_md:
        return native_md, "pypdf-fallback", metadata
    raise RuntimeError(
        f"No parser produced content for {path.name}; classification={classification}, "
        f"OCR_PROVIDER={ocr_config.get('provider') or OCR_PROVIDER}, "
        f"LOCAL_DOCLING_ENABLED={LOCAL_DOCLING_ENABLED}"
    )


def assess_content_quality(markdown: str, suffix: str, task: dict[str, Any]) -> dict[str, Any]:
    """Return a conservative quality signal before a document is published.

    The parser may produce non-empty but unusable output (font encoding
    damage, a blank image, or a PPTX with only unrecognised shapes). Quality is
    therefore persisted separately from parser success. A review result keeps
    the Markdown/chunks available for inspection but never enters GBrain.
    """
    text = str(markdown or "")
    placeholder_count = len(re.findall(r"<!--\s*(?:image|picture|figure)\s*-->", text, flags=re.IGNORECASE))
    quality_text = re.sub(r"<!--.*?-->", "", text, flags=re.DOTALL)
    quality_text = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", quality_text)
    meaningful = re.findall(r"[A-Za-z0-9\u4e00-\u9fff]", quality_text)
    meaningful_count = len(meaningful)
    replacement_count = text.count("\ufffd")
    control_count = sum(1 for char in text if ord(char) < 32 and char not in "\n\r\t")
    replacement_ratio = replacement_count / max(len(text), 1)
    control_ratio = control_count / max(len(text), 1)
    issues: list[str] = []
    binary_input = suffix not in {".md", ".txt", ".csv", ".html", ".htm"}
    if meaningful_count == 0:
        issues.append("没有提取到可检索文字")
    if binary_input and meaningful_count < 20:
        issues.append("提取文字过少，可能是空白文件或解析不完整")
    if replacement_ratio > 0.01:
        issues.append("存在较多字体编码替换字符")
    if control_ratio > 0.02:
        issues.append("存在异常控制字符")
    if placeholder_count and suffix in {".pptx", ".png", ".jpg", ".jpeg"}:
        issues.append("版面解析只返回图片占位符，图片文字尚未完成 OCR")
    confidence = task.get("ocr_average_confidence")
    if confidence is not None:
        try:
            if float(confidence) < 0.75:
                issues.append("OCR 平均置信度低于 0.75")
        except (TypeError, ValueError):
            pass
    score = 1.0
    score -= min(replacement_ratio * 4, 0.45)
    score -= min(control_ratio * 2, 0.2)
    if binary_input and meaningful_count < 20:
        score -= 0.45
    if confidence is not None:
        try:
            score = min(score, max(0.0, float(confidence)))
        except (TypeError, ValueError):
            pass
    score = round(max(0.0, min(1.0, score)), 4)
    status = "passed" if not issues else "needs_review"
    if meaningful_count == 0:
        status = "rejected"
    return {
        "quality_status": status,
        "quality_score": score,
        "quality_issues": issues,
        "quality_metrics": {
            "characters": len(text),
            "meaningful_characters": meaningful_count,
            "replacement_ratio": round(replacement_ratio, 6),
            "control_ratio": round(control_ratio, 6),
            "image_placeholders": placeholder_count,
        },
    }


async def process_file(
    task_id: str,
    path: Path,
    parser_type: str,
    ocr_config: dict[str, str],
) -> None:
    task = tasks[task_id]
    task["status"] = "processing"
    try:
        suffix = path.suffix.lower()
        if suffix in {".md", ".txt", ".csv", ".html", ".htm"}:
            content = await asyncio.to_thread(path.read_bytes)
            task["markdown"] = extract_plaintext(path.name, content)
            task["engine"] = "plaintext"
        elif suffix == ".doc":
            task["conversion"] = "antiword"
            task["markdown"] = await asyncio.to_thread(extract_legacy_word, path)
            task["engine"] = "antiword"
        elif suffix == ".docx":
            md = await asyncio.to_thread(extract_docx, path)
            if md:
                task["markdown"] = md
                task["engine"] = "python-docx"
            else:
                if not LOCAL_DOCLING_ENABLED:
                    raise RuntimeError("DOCX native extraction returned no text and local Docling is disabled")
                md = await convert_with_docling(path)
                task["markdown"] = md
                task["engine"] = "docling"
        elif suffix in {".xlsx", ".xls"}:
            md = await asyncio.to_thread(extract_excel, path)
            if md:
                task["markdown"] = md
                task["engine"] = "openpyxl" if suffix == ".xlsx" else "xlrd"
            else:
                if not LOCAL_DOCLING_ENABLED:
                    raise RuntimeError("Excel native extraction returned no cells and local Docling is disabled")
                md = await convert_with_docling(path)
                task["markdown"] = md
                task["engine"] = "docling"
        elif suffix == ".pptx":
            if LOCAL_DOCLING_ENABLED:
                try:
                    md = await convert_with_docling(path)
                    if re.search(r"<!--\s*(?:image|picture|figure)\s*-->", md, flags=re.IGNORECASE):
                        raise RuntimeError("Docling returned image placeholders; OCR is required for complete PPTX ingestion")
                    task["markdown"] = md
                    task["engine"] = "docling-local"
                except Exception as docling_error:
                    logger.warning("Local Docling PPTX conversion failed for %s: %s", path.name, docling_error)
                    task["docling_error"] = str(docling_error)
                    md, engine, parser_metadata = await convert_pptx_without_docling(path, ocr_config)
                    task["markdown"] = md
                    task["engine"] = engine
                    task.update(parser_metadata)
            else:
                md, engine, parser_metadata = await convert_pptx_without_docling(path, ocr_config)
                task["markdown"] = md
                task["engine"] = engine
                task.update(parser_metadata)
        elif suffix == ".pdf":
            pdf_info = await asyncio.to_thread(inspect_pdf_native, path)
            native_md = str(pdf_info.get("markdown", ""))
            for key in (
                "classification",
                "page_count",
                "text_pages",
                "native_chars",
                "native_page_ratio",
                "native_quality",
            ):
                task[key] = pdf_info.get(key)
            md, engine, parser_metadata = await convert_pdf_with_fallback(
                path,
                str(pdf_info.get("classification") or "unknown"),
                native_md,
                ocr_config,
                [str(text) for text in pdf_info.get("page_texts", [])],
                [int(index) for index in pdf_info.get("native_page_indexes", [])],
            )
            task["markdown"] = md
            task["engine"] = engine
            task.update(parser_metadata)
        else:
            # Standalone images use local Docling in the test profile and the
            # cloud OCR route in production. No image is silently accepted
            # without text extraction.
            if LOCAL_DOCLING_ENABLED:
                try:
                    md = await convert_with_docling(path)
                    if re.search(r"<!--\s*(?:image|picture|figure)\s*-->", md, flags=re.IGNORECASE):
                        raise RuntimeError("Docling returned an image placeholder; OCR is required for complete image ingestion")
                    task["markdown"] = md
                    task["engine"] = "docling-local"
                except Exception as docling_error:
                    logger.warning("Local Docling image conversion failed for %s: %s", path.name, docling_error)
                    task["docling_error"] = str(docling_error)
                    md, ocr_metadata = await convert_image_with_baidu_ocr(path, ocr_config)
                    task["markdown"] = f"# {path.stem}\n\n{md}"
                    task["engine"] = "ocr-baidu-image"
                    task.update(ocr_metadata)
            else:
                md, ocr_metadata = await convert_image_with_baidu_ocr(path, ocr_config)
                task["markdown"] = f"# {path.stem}\n\n{md}"
                task["engine"] = "ocr-baidu-image"
                task.update(ocr_metadata)

        if not task.get("markdown", "").strip():
            raise RuntimeError("Extracted Markdown is empty")
        task["markdown"] = normalize_markdown(
            task["markdown"].replace("\x00", "").replace("\u0000", ""),
            str(task.get("filename", "upload.md")),
        )
        task.update(assess_content_quality(task["markdown"], suffix, task))
        if task["quality_status"] == "rejected":
            raise RuntimeError("Quality gate rejected the document: " + "; ".join(task["quality_issues"]))
        task["status"] = "completed"
        logger.info(
            "Task %s completed successfully via engine=%s classification=%s "
            "ocr_provider=%s native_page_ratio=%s",
            task_id,
            task.get("engine"),
            task.get("classification", "n/a"),
            task.get("ocr_provider", "none"),
            task.get("native_page_ratio", "n/a"),
        )
    except Exception as exc:
        task["status"] = "failed"
        task["error"] = f"Processing failed for {task.get('filename', 'unknown')}: {type(exc).__name__} ({exc})"
        logger.error(f"Error processing file for task {task_id}: {exc}", exc_info=True)
    finally:
        try:
            path.unlink(missing_ok=True)
        except Exception:
            pass

@app.post("/parse", response_model=ParseResponse)
async def parse_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    parser_type: str = "docling",
    ocr_provider: str | None = Form(None),
    ocr_endpoint: str | None = Form(None),
    ocr_api_key: str | None = Form(None),
    ocr_secret_key: str | None = Form(None),
    _auth: None = Depends(verify_auth),
):
    if len(tasks) >= 5000:
        oldest_tid = min(tasks.keys(), key=lambda k: tasks[k].get("created_at", float('inf')))
        del tasks[oldest_tid]

    filename = Path(file.filename or "upload.md").name
    suffix = Path(filename).suffix.lower()
    if not suffix:
        suffix = ".md"
        filename = f"{filename}.md"
    if suffix not in SUPPORTED_EXTENSIONS:
        raise HTTPException(status_code=415, detail=f"Unsupported file type: {suffix or 'unknown'}")
        
    if file.size and file.size > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail="File exceeds 200 MiB limit")
        
    content = await file.read(MAX_FILE_BYTES + 1)
    if len(content) > MAX_FILE_BYTES:
        raise HTTPException(status_code=413, detail="File exceeds 200 MiB limit")

    task_id = str(uuid.uuid4())
    path = UPLOAD_ROOT / f"{task_id}{suffix}"
    await asyncio.to_thread(path.write_bytes, content)
    tasks[task_id] = {"status": "queued", "filename": filename, "parser_type": parser_type, "created_at": time.time()}
    # Credentials are request-scoped and deliberately not copied into tasks;
    # /parse/{task_id} must never expose them.
    ocr_config = {
        "provider": (ocr_provider or OCR_PROVIDER).strip().lower(),
        "endpoint": (ocr_endpoint or BAIDU_OCR_ENDPOINT).strip(),
        "api_key": ocr_api_key or BAIDU_OCR_API_KEY,
        "secret_key": ocr_secret_key or BAIDU_OCR_SECRET_KEY,
    }
    background_tasks.add_task(process_file, task_id, path, parser_type, ocr_config)
    return ParseResponse(task_id=task_id, status="accepted", message=f"File {filename} queued for parsing")

@app.get("/parse/{task_id}")
def parse_status(task_id: str):
    task = tasks.get(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="Parse task not found")
    return {"task_id": task_id, **task}

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8100)
