from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
IMAGE_DIR = ROOT / "design" / "government-introduction" / "images"
OUTPUT = ROOT / "design" / "government-introduction" / "政务知识智能中枢-政府介绍材料-V2.0.pptx"


def main() -> None:
    images = sorted(IMAGE_DIR.glob("*.png"))
    if len(images) != 8:
        raise SystemExit(f"expected 8 slide images, found {len(images)}")

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image), 0, 0, width=prs.slide_width, height=prs.slide_height
        )

    prs.core_properties.title = "政务知识智能中枢——政府介绍材料"
    prs.core_properties.subject = "企业级知识库与可信问答平台"
    prs.core_properties.author = "GBrain Knowledge Platform"
    prs.core_properties.comments = "基于项目当前权限管理、入库、检索和局域网部署能力制作"
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
